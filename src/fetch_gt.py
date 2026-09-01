"""GT 建库 —— 解析通道（非 HTML 文档）。

docfmt 那 22 条是全场 GT 质量最好的一批：CSV 的 GT 就是那个 CSV，xlsx 的 GT 就是解析出的表。
不需要浏览器、不需要 key，纯本地。

**本轮只评抓取能力**，所以解析器只负责把文档变成可比对的文本（供 coverage 用），
不再产出结构计数 —— 表格/幻灯片保真属于解析质量，已从评价里删除。

解析失败一律返回规则名，不抛 —— 调用方要能区分"这页解析不了"和"我们的解析器崩了"。
"""
from __future__ import annotations

import signal
import threading
from contextlib import contextmanager

import csv as _csv
import io
import json as _json
import re

_MAGIC = ((b"%PDF", "pdf"), (b"PK\x03\x04", "zip"))
_ZIP_KINDS = {"docx", "xlsx", "pptx"}
_CT_MAP = {
    "application/pdf": "pdf",
    "text/csv": "csv",
    "application/json": "json",
    "application/rss+xml": "rss",
    "application/atom+xml": "atom",
    "text/plain": "txt",
    "text/markdown": "md",
    "application/xml": "xml",
    "text/xml": "xml",
    "text/html": "html",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}


def _import(name: str):
    """单点 import，便于测试打桩成"依赖缺失"。"""
    try:
        return __import__(name)
    except ImportError:
        return None


def sniff_doc_type(raw: bytes, content_type: str, declared: str) -> tuple[str, str]:
    """(doc_type, 规则名)。三级：魔数 > content-type > URL 声明。

    `arxiv.org/pdf/1706.03762` 没有后缀，`declared` 是 unknown，靠后两级定 —— 这道题
    考的就是各家会不会认错。规则名逐条记录，事后能核是哪一级判出来的（playbook §9.3）。
    """
    head = (raw or b"")[:8]
    for sig, kind in _MAGIC:
        if head.startswith(sig):
            if kind == "zip":
                # docx/xlsx/pptx 都是 zip，光看魔数分不出来
                if declared in _ZIP_KINDS:
                    return declared, "declared_zip"
                return "zip", "magic_bytes"
            return kind, "magic_bytes"
    ct = (content_type or "").split(";")[0].strip().lower()
    if ct in _CT_MAP:
        return _CT_MAP[ct], "content_type"
    if declared and declared != "unknown":
        return declared, "declared"
    return "unknown", "undetermined"


def _table_md(rows: list[list]) -> tuple[str, int]:
    """表格渲染成 Markdown 管道表，返回 (文本, 数据行数)。

    表头 + 分隔行 + 数据行 —— checks 侧数管道行再减 2，两边对得上。
    """
    if not rows:
        return "", 0
    head = rows[0]
    out = ["| " + " | ".join(str(c if c is not None else "") for c in head) + " |",
           "| " + " | ".join("---" for _ in head) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(str(c if c is not None else "") for c in r) + " |")
    return "\n".join(out), max(0, len(rows) - 1)


def _ok(text: str, doc_type: str, rule: str = "parsed") -> dict:
    if not (text or "").strip():
        rule = "parsed_no_text"
    return {"text": text or "", "doc_type": doc_type, "rule": rule}


def _fail(doc_type: str, rule: str) -> dict:
    return {"text": "", "doc_type": doc_type, "rule": rule}


def parse_document(raw: bytes, doc_type: str, url: str) -> dict:
    """把一份非 HTML 文档解析成 (Markdown 文本 + 结构计数)。**失败返回规则名，不抛。**"""
    try:
        return _PARSERS.get(doc_type, _parse_text)(raw, doc_type)
    except Exception:                            # noqa: BLE001
        return _fail(doc_type, "parse_failed")


def _parse_csv(raw: bytes, doc_type: str) -> dict:
    rows = list(_csv.reader(io.StringIO(raw.decode("utf-8", "replace"))))
    rows = [r for r in rows if any(str(c).strip() for c in r)]
    text, n = _table_md(rows)
    return _ok(text, doc_type)


def _parse_json(raw: bytes, doc_type: str) -> dict:
    obj = _json.loads(raw.decode("utf-8", "replace"))
    text = _json.dumps(obj, ensure_ascii=False, indent=2)
    return _ok(text, doc_type)


def _parse_feed(raw: bytes, doc_type: str) -> dict:
    fp = _import("feedparser")
    if fp is None:
        return _fail(doc_type, "dep_missing")
    d = fp.parse(raw)
    items = d.get("entries") or []
    text = "\n".join("- " + (e.get("title") or "").strip() for e in items)
    return _ok(text, doc_type)


def _parse_xml(raw: bytes, doc_type: str) -> dict:
    s = raw.decode("utf-8", "replace")
    nodes = len(re.findall(r"<(?!\?|!)([a-zA-Z][\w:-]*)", s))
    text = re.sub(r"<[^>]+>", " ", s)
    text = re.sub(r"\s+", " ", text).strip()
    return _ok(text, doc_type)


def _parse_pdf(raw: bytes, doc_type: str) -> dict:
    pypdf = _import("pypdf")
    if pypdf is None:
        return _fail(doc_type, "dep_missing")
    r = pypdf.PdfReader(io.BytesIO(raw))
    parts = []
    for i, page in enumerate(r.pages, 1):
        parts.append("## Page %d" % i)
        parts.append((page.extract_text() or "").strip())
    return _ok("\n\n".join(p for p in parts if p), doc_type)


def _parse_docx(raw: bytes, doc_type: str) -> dict:
    if _import("docx") is None:
        return _fail(doc_type, "dep_missing")
    import docx
    d = docx.Document(io.BytesIO(raw))
    paras = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    blocks = list(paras)
    for t in d.tables:
        rows = [[c.text for c in row.cells] for row in t.rows]
        md, _ = _table_md(rows)
        blocks.append(md)
    return _ok("\n\n".join(blocks),
               {"paragraphs": len(paras), "tables": len(d.tables)}, doc_type)


def _parse_xlsx(raw: bytes, doc_type: str) -> dict:
    if _import("openpyxl") is None:
        return _fail(doc_type, "dep_missing")
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True)
    blocks, total = [], 0
    for ws in wb.worksheets:
        rows = [list(r) for r in ws.iter_rows(values_only=True)
                if any(c is not None and str(c).strip() for c in r)]
        md, n = _table_md(rows)
        total += n
        # sheets 只进明细表不进分数 —— 厂商的 markdown 不会带这个标记
        blocks.append("## Sheet: %s\n%s" % (ws.title, md))
    return _ok("\n\n".join(blocks),
               {"sheets": len(wb.worksheets), "rows": total}, doc_type)


def _parse_pptx(raw: bytes, doc_type: str) -> dict:
    if _import("pptx") is None:
        return _fail(doc_type, "dep_missing")
    from pptx import Presentation
    pres = Presentation(io.BytesIO(raw))
    blocks = []
    for i, slide in enumerate(pres.slides, 1):
        texts = [sh.text.strip() for sh in slide.shapes
                 if getattr(sh, "has_text_frame", False) and sh.text.strip()]
        blocks.append("## Slide %d\n%s" % (i, "\n".join(texts)))
    return _ok("\n\n".join(blocks), doc_type)


def _parse_text(raw: bytes, doc_type: str) -> dict:
    s = raw.decode("utf-8", "replace")
    lines = [l for l in s.splitlines() if l.strip()]
    return _ok(s, doc_type)


_PARSERS = {
    "csv": _parse_csv, "json": _parse_json,
    "rss": _parse_feed, "atom": _parse_feed,
    "xml": _parse_xml, "pdf": _parse_pdf,
    "docx": _parse_docx, "xlsx": _parse_xlsx, "pptx": _parse_pptx,
    "txt": _parse_text, "md": _parse_text,
}


# ══════════════════════════════════════════════════════════════════════════
# 浏览器通道：HTML 页的 GT
# ══════════════════════════════════════════════════════════════════════════

from collections import Counter                                    # noqa: E402

from .fetch_checks import tokenize                                 # noqa: E402
from .fetch_spec import TH                                         # noqa: E402

# 排除这些子树之后剩下的是正文；被排除的那部分恰好就是样板词的来源。
# **精度口径的分母由此而来** —— 比链接密度启发式硬得多。
_CHANNELS = frozenset({"playwright_headless", "chrome_real"})

BOILER_SELECTORS = ("nav", "header", "footer", "aside", "[role=navigation]",
                    "[role=banner]", "[role=contentinfo]", ".nav", ".navbar",
                    ".header", ".footer", ".sidebar", ".menu", ".breadcrumb",
                    "#nav", "#header", "#footer", "#sidebar")

VOCAB_CAP = 60
_STOP = frozenset("""
a an the and or but if then else of in on at to for from by with without into onto
is are was were be been being do does did have has had will would can could should
this that these those it its as not no yes you your we our they their he she his her
i me my us them there here what which who whom when where why how all any both each
more most other some such only own same so than too very s t don now
""".split())


def _playwright():
    """单点 import，测试据此跳过。"""
    try:
        from playwright.sync_api import sync_playwright
        return sync_playwright
    except ImportError:
        return None


def _content_terms(text: str, exclude: set[str], cap: int) -> list[str]:
    """按词频取内容词：去停用词、去单字符拉丁、去 exclude 集合。"""
    freq = Counter(t for t in tokenize(text)
                   if t not in _STOP and t not in exclude
                   and (len(t) >= 2 or not t.isascii()))
    return [w for w, _ in freq.most_common(cap)]


def derive_vocab(main_text: str, boiler_text: str) -> dict:
    """产出内容词表（成功闸门的分母）。样板词表只作排除集，不是指标。

    尾部词取正文末 10%，喂 `oversize` 的截断检查 —— 静默截断在 coverage 口径下照样是
    pass，只要词表落在前半。
    """
    # **样板词的判据是"在样板侧比在正文侧更常见"，不是"在样板侧出现过"。**
    # docs.python.org 的侧边栏目录列出了全部章节名，而那些恰恰是这页最核心的内容词
    # （task / coroutines / cancellation）。只要"出现过"就排除，召回的分母就丢了
    # 最该考的那些词 —— 比 noise_ratio 变糊严重得多。
    main_freq = Counter(t for t in tokenize(main_text) if t not in _STOP)
    boiler_freq = Counter(t for t in tokenize(boiler_text) if t not in _STOP)
    boiler_only = {t for t, n in boiler_freq.items() if n >= main_freq.get(t, 0)}
    boiler = _content_terms(boiler_text, set(boiler_freq) - boiler_only, VOCAB_CAP)
    vocab = _content_terms(main_text, boiler_only, VOCAB_CAP)
    toks = tokenize(main_text)
    # **身份锚点要从文档开头取**：全局罕见词会挑到表格深处的值（arxiv 那篇论文挑出了
    # wsj / 38 / pad），那是**完整度**标记不是**身份**标记 —— 抽到了正确的论文但表格
    # 没抽全的家会被判成"返回错页"。"这是不是这一页"看开头就能答。
    head_src = toks[:max(200, int(len(toks) * 0.1))]
    head = _content_terms(" ".join(head_src), boiler_only, VOCAB_CAP)
    return {"vocab": vocab, "vocab_n": len(vocab),
            "boiler_terms": boiler, "vocab_head": head,
            "degenerate": len(vocab) < TH["vocab_min"]}


def gt_is_walled(main_text: str, title: str) -> list[str]:
    """GT 自己是不是一堵墙。

    浏览器架构上就是 Firecrawl / Zyte 卖的东西，防守页它也会吃拦截。渲染"成功"且有
    词表并不等于拿到了内容 —— w3.org/TR/html52/ 实测建出来的 GT 标题是 "Just a moment..."，
    整页词表是 Cloudflare 验证页的文案，于是全场五家都在跟一张验证页比对。
    只在"词表为空"时标缺口兜不住这一类（设计文档 二.4）。
    """
    from .fetch_checks import wall_hit
    return sorted(set(wall_hit(main_text) + wall_hit(title or "")))


def derive_strength(headless_ok: bool | None, chrome_ok: bool | None) -> str:
    """防护强度是两条 GT 通道的免费副产品（设计文档 §1.4）。

    **两条通道各自"没跑过"都要落到 unknown，不能当成"跑了但失败"。**
    默认 hard 会把"没测"伪装成"测出来最难"；而把缺失的 headless 当成 False，会让
    "只跑过真 Chrome 且成功"的页判成 medium —— 那同样是凭空造出来的结论。
    """
    if headless_ok:
        return "soft"
    if headless_ok is None or chrome_ok is None:
        return "unknown"
    return "medium" if chrome_ok else "hard"


_JS_EXTRACT = """
() => {
  const sels = %s;
  const doc = document;
  doc.querySelectorAll('script,style,noscript,template').forEach(n => n.remove());
  const boilerNodes = [];
  sels.forEach(s => { try { doc.querySelectorAll(s).forEach(n => boilerNodes.push(n)); }
                      catch (e) {} });
  const boilerText = boilerNodes.map(n => n.innerText || '').join('\\n');
  // 临时隐藏样板节点，在**活文档**上取 innerText —— 克隆节点没有布局，
  // innerText 会退化成 textContent，把 <style> 里的 CSS 也算成正文。
  const prev = boilerNodes.map(n => n.style.display);
  boilerNodes.forEach(n => { n.style.display = 'none'; });
  const mainText = doc.body.innerText || '';
  boilerNodes.forEach((n, i) => { n.style.display = prev[i]; });
  return {mainText, boilerText, title: doc.title || ''};
}
"""


def render_page(url: str, *, channel: str = "playwright_headless",
                timeout: int = 45, screenshot_dir: str | None = None,
                pid: str = "page") -> dict:
    """渲染一页并取 GT。**干净 profile、不登录、只接 cookie。**

    用已登录会话建 GT，登录墙页的 GT 就包含任何 provider 都拿不到的内容，
    然后全场对着一份谁都达不到的答案一起失败 —— 那不是在量 provider（设计文档 §〇.6）。
    """
    sync_playwright = _playwright()
    if sync_playwright is None:
        return {"error": "playwright 未安装", "rule": "dep_missing"}
    if channel not in _CHANNELS:
        raise ValueError("未知通道 %r，可选 %s" % (channel, sorted(_CHANNELS)))
    budget = timeout * 2 + 30          # 导航 + 滚动 + 截图，留足再加一倍余量
    try:
        with _watchdog(budget):
            got = _render_once(sync_playwright, url, channel, timeout,
                               screenshot_dir, pid, [])
    except TimeoutError as e:
        return {"channel": channel, "url": url, "rule": "render_failed",
                "error": str(e), "main_text": "", "boiler_text": ""}
    if got.get("rule") == "render_failed" and "ERR_HTTP2" in (got.get("error") or ""):
        # 有些站的 HTTP/2 实现和 chromium 谈不拢（bestbuy / noon 实测）。关掉 H2 重试
        # 一次再下结论 —— 不重试的话"协议谈不拢"会被记成"这页抓不到"。
        try:
            with _watchdog(budget):
                got = _render_once(sync_playwright, url, channel, timeout,
                                   screenshot_dir, pid, ["--disable-http2"])
        except TimeoutError as e:
            return {"channel": channel, "url": url, "rule": "render_failed",
                    "error": str(e), "main_text": "", "boiler_text": ""}
        got["retried_without_http2"] = True
    return got


@contextmanager
def _watchdog(seconds: int):
    """硬看门狗。**Playwright 的 timeout 只管单个操作**，页面在浏览器进程里挂住时
    整轮会无限期堵死 —— 真 Chrome 跑反爬页时实测卡了 20 分钟一行没落盘。
    只在主线程生效（SIGALRM 的限制），浏览器通道本来就是顺序跑的。
    """
    if seconds <= 0 or threading.current_thread() is not threading.main_thread():
        yield
        return

    def _boom(signum, frame):
        raise TimeoutError("render 看门狗触发：单页超过 %ds" % seconds)

    old = signal.signal(signal.SIGALRM, _boom)
    signal.alarm(seconds)
    try:
        yield
    finally:
        signal.alarm(0)
        signal.signal(signal.SIGALRM, old)


def _render_once(sync_playwright, url: str, channel: str, timeout: int,
                 screenshot_dir: str | None, pid: str, args: list[str]) -> dict:
    import json as _j
    out: dict = {"channel": channel, "url": url}
    with sync_playwright() as p:
        # **channel 必须真的改变启动方式。** 之前这里写死 headless=True，`chrome_real`
        # 只是个传着好看的参数 —— 那条命令会静默跑成 headless 却声称用了真 Chrome，
        # 而反爬页的整个意义就在于真实浏览器指纹。
        if channel == "chrome_real":
            # 本机装的 Google Chrome，有头、真实指纹、**干净 profile 不登录**
            # （用登录态建 GT 等于让全场对着一份谁都拿不到的答案一起失败）
            browser = p.chromium.launch(headless=False, channel="chrome", args=args)
        else:
            browser = p.chromium.launch(headless=True, args=args)
        ctx = browser.new_context(
            user_agent=("Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                        "AppleWebKit/537.36 (KHTML, like Gecko) "
                        "Chrome/125.0 Safari/537.36"),
            viewport={"width": 1280, "height": 900}, locale="en-US")
        page = ctx.new_page()
        try:
            resp = page.goto(url, timeout=timeout * 1000, wait_until="domcontentloaded")
            out["http_status"] = resp.status if resp else None
            _scroll_to_bottom(page)
            data = page.evaluate(_JS_EXTRACT % _j.dumps(list(BOILER_SELECTORS)))
            out["main_text"] = data["mainText"]
            out["boiler_text"] = data["boilerText"]
            out["title"] = data["title"]
            out["rule"] = "rendered"
            if screenshot_dir:
                # **截图单独兜底。** 截图失败不该拖垮已经取到的正文 —— hub.docker.com
                # 实测 http=200、正文正常，只是 captureScreenshot 报协议错，
                # 结果整页 GT 被丢掉了。
                try:
                    import os as _os
                    _os.makedirs(screenshot_dir, exist_ok=True)
                    shot = _os.path.join(screenshot_dir, "%s.png" % pid)
                    page.screenshot(path=shot, full_page=True)
                    out["screenshot_path"] = shot
                except Exception as e:           # noqa: BLE001
                    out["screenshot_error"] = "%s: %s" % (type(e).__name__, str(e)[:120])
        except Exception as e:                   # noqa: BLE001
            out.update({"error": "%s: %s" % (type(e).__name__, str(e)[:200]),
                        "rule": "render_failed", "main_text": "", "boiler_text": ""})
        finally:
            for closer in (ctx.close, browser.close):
                try:
                    closer()
                except Exception:                # noqa: BLE001  看门狗触发时 close 也可能挂
                    pass
    return out


def _scroll_to_bottom(page, rounds: int = 3, quiet_ms: int = 2000) -> None:
    """滚到底触发懒加载。**有确定的停止条件**：高度不再增长或轮数用尽。"""
    last = 0
    for _ in range(rounds):
        page.mouse.wheel(0, 20000)
        try:
            page.wait_for_load_state("networkidle", timeout=quiet_ms)
        except Exception:                        # noqa: BLE001
            pass
        h = page.evaluate("() => document.body.scrollHeight")
        if h == last:
            break
        last = h


def derive_anchors(title: str, vocab: list[str], df: dict[str, int],
                   n_pages: int, cap: int = 6, url: str = "") -> list[str]:
    """页面独有锚点：标题实词 + 在全集里罕见的内容词。

    只有 100 页语料，"独有"判不准，所以用"在这 100 页里出现在少于 20% 的页上"当近似。
    这一处已知不完美，写进了设计文档的开放项。
    """
    # 文档频率门限收到 5%。两成在 58 页语料上是 11 页 —— 而集合里 PDF 才 8 个，
    # "pdf" 就这么混进了锚点。锚点虚高不会造出假 lost（方向是安全的），但会让
    # "返回错页"这条硬否决形同虚设：返回另一份 IRS 表格照样命中 line/see/tax。
    lim = max(1, int(n_pages * 0.05))
    # URL 末段是最强的身份信号（f1040 / rfc2616 / asyncio-task），优先于标题与词表
    slug_src = re.sub(r"[^a-z0-9]+", " ", (url or "").lower().rsplit("/", 2)[-1])
    out = [t for t in _content_terms(slug_src, set(), 4) if df.get(t, 0) <= lim]
    for t in _content_terms(title or "", set(), 6):
        if t not in out and df.get(t, 0) <= lim:
            out.append(t)
    rare = [t for t in vocab if t not in _STOP and df.get(t, 0) <= lim]
    for t in rare:
        if t not in out:
            out.append(t)
        if len(out) >= cap:
            break
    return out[:cap]
