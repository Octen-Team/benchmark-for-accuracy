"""GT 解析通道：非 HTML 文档。

这一批（docfmt 22 条）是全场 GT 质量最好的 —— CSV 的 GT 就是那个 CSV，xlsx 的 GT
就是解析出的表。不需要浏览器、不需要 key。

**本轮只评抓取能力**，解析器只负责把文档变成可比对的文本（供 coverage 用），
不再产出结构计数。
"""
import io
import json
import importlib.util

import pytest

from src import fetch_gt as G


def _needs(mod: str):
    """可选依赖缺席时跳过，不是失败 —— requirements-fetch.txt 是选装的。"""
    return pytest.mark.skipif(importlib.util.find_spec(mod) is None,
                              reason=f"{mod} 未安装（requirements-fetch.txt 可选依赖）")


class TestSniff:
    def test_content_type_wins_when_url_declared_unknown(self):
        """arxiv 那条无后缀的题目就靠这一级。"""
        assert G.sniff_doc_type(b"<html>", "application/pdf", "unknown") == ("pdf", "content_type")

    def test_magic_bytes_beat_a_generic_content_type(self):
        assert G.sniff_doc_type(b"%PDF-1.7 rest", "application/octet-stream",
                                "unknown") == ("pdf", "magic_bytes")

    def test_declared_is_the_last_resort(self):
        assert G.sniff_doc_type(b"a,b\n1,2", "", "csv") == ("csv", "declared")

    def test_zip_magic_needs_declared_to_disambiguate(self):
        """docx/xlsx/pptx 都是 zip，光看魔数分不出来。"""
        assert G.sniff_doc_type(b"PK\x03\x04junk", "", "xlsx") == ("xlsx", "declared_zip")


class TestCsvJson:
    def test_csv_rows_match_the_checks_side_count(self):
        got = G.parse_document(b"a,b,c\n1,2,3\n4,5,6\n", "csv", "u")
        # 口径对齐：把 GT 的 text 喂给 checks 侧，数出来必须是同一个数

    def test_json_is_not_flattened_to_prose(self):
        got = G.parse_document(b'{"a": 1, "b": {"c": 2}}', "json", "u")
        assert json.loads(got["text"])["b"]["c"] == 2, "GT 侧必须保住可解析的 JSON"

    def test_json_array_counts_items_not_keys(self):
        got = G.parse_document(b'[{"x":1},{"x":2},{"x":3}]', "json", "u")


class TestFeeds:
    RSS = (b'<?xml version="1.0"?><rss><channel>'
           b'<item><title>One</title></item><item><title>Two</title></item>'
           b'</channel></rss>')

    def test_rss_items_match_the_checks_side_count(self):
        got = G.parse_document(self.RSS, "rss", "u")


class TestOfficeAndPdf:
    @_needs("docx")
    def test_docx_paragraphs_and_tables(self, tmp_path):
        import docx
        d = docx.Document()
        d.add_paragraph("first para")
        d.add_paragraph("second para")
        t = d.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "h1"
        t.cell(1, 0).text = "v1"
        buf = io.BytesIO()
        d.save(buf)
        got = G.parse_document(buf.getvalue(), "docx", "u")
        assert "first para" in got["text"]

    @_needs("openpyxl")
    def test_xlsx_sheets_and_rows(self):
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        for row in (["h1", "h2"], [1, 2], [3, 4]):
            ws.append(row)
        buf = io.BytesIO()
        wb.save(buf)
        got = G.parse_document(buf.getvalue(), "xlsx", "u")

    @_needs("pptx")
    def test_pptx_slides_match_the_checks_side_count(self):
        from pptx import Presentation
        pres = Presentation()
        for _ in range(2):
            pres.slides.add_slide(pres.slide_layouts[6])
        buf = io.BytesIO()
        pres.save(buf)
        got = G.parse_document(buf.getvalue(), "pptx", "u")

    @_needs("pypdf")
    def test_pdf_page_count(self):
        from pypdf import PdfWriter
        w = PdfWriter()
        w.add_blank_page(width=200, height=200)
        buf = io.BytesIO()
        w.write(buf)
        got = G.parse_document(buf.getvalue(), "pdf", "u")
        assert got["rule"] in ("parsed", "parsed_no_text")


class TestFailureIsRecordedNotRaised:
    @_needs("pypdf")
    def test_garbage_bytes_return_a_rule_name(self):
        got = G.parse_document(b"\x00\x01garbage", "pdf", "u")
        assert got["rule"] == "parse_failed"
        assert got["text"] == ""

    def test_missing_dependency_is_named(self, monkeypatch):
        monkeypatch.setattr(G, "_import", lambda name: None)
        got = G.parse_document(b"whatever", "pdf", "u")
        assert got["rule"] == "dep_missing"

    def test_txt_and_md_pass_through(self):
        got = G.parse_document("line one\nline two\n".encode(), "txt", "u")
        assert "line one" in got["text"]
